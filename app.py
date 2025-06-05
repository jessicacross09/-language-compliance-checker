import streamlit as st
import os
import re
import pandas as pd
from io import StringIO
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
import spacy

# Load spaCy model
nlp = spacy.load("en_core_web_sm")

# --- Streamlit Config ---
st.set_page_config(page_title="APEC-RISE Text Harmonization Tool", layout="wide")
st.title("\U0001F50D APEC-RISE Text Harmonization Tool")
st.caption("Scan documents for non-compliant terms and download flagged text with highlights.")
st.divider()

# --- Banned Terms Dictionary ---
banned_terms_dict = {
    "accessible": ["reachable; usable; broadly available"],
    "activism": ["stakeholder engagement; issue-focused participation"],
    "anti-racism": ["addressing discrimination"],
    "bias": ["subjective influence"],
    "clean energy": ["energy diversification; secure energy sourcing"],
    "climate change": ["environmental shifts"],
    "climate crisis": ["environmental conditions; weather-related impacts"],
    "climate science": ["environmental data"],
    "country": ["economy"],
    "DEI": ["broad participation; inclusive stakeholder engagement"],
    "DEIA": ["broad participation; inclusive stakeholder engagement"],
    "disability": ["persons with disabilities"],
    "diverse": ["varied"],
    "diverse backgrounds": ["varied experiences"],
    "diverse communities": ["different population groups"],
    "diverse community": ["broad-based community"],
    "diverse group": ["multifaceted group"],
    "diverse groups": ["multiple stakeholder groups"],
    "diversified": ["varied"],
    "diversify": ["broaden"],
    "diversifying": ["expanding"],
    "diversity": ["variety"],
    "equal opportunity": ["merit-based opportunity; fair process"],
    "equity": ["fair access"],
    "gender": ["demographics"],
    "gender mainstreaming": ["gender considerations"],
    "gender-responsive": ["inclusive of gender perspectives"],
    "Gulf of Mexico": ["Gulf of America"],
    "inclusion": ["broad-based participation"],
    "inclusive": ["participatory"],
    "inclusive leadership": ["broad-based leadership"],
    "inclusiveness": ["broad engagement"],
    "inclusivity": ["collaborative approaches"],
    "inequality": ["gaps in access or opportunity"],
    "national": ["domestic"],
    "non-binary": ["gender-diverse"],
    "nonbinary": ["gender-diverse"],
    "oppression": ["restrictive conditions; limiting environments"],
    "oppressive": ["restrictive; limiting"],
    "social justice": ["fair policy outcomes; community fairness"],
    "socioeconomic": ["economic background; living conditions"],
    "Taiwan": ["Chinese Taipei"],
    "victim": ["impacted individual; affected party"],
    "vulnerable populations": ["underrepresented stakeholders"]
}

# --- Named Entity Exception ---
def is_named_entity(snippet: str, term: str):
    doc = nlp(snippet)
    for ent in doc.ents:
        if term.lower() in ent.text.lower() and ent.label_ == "ORG":
            return True
    return False

# --- File Readers ---
def read_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def read_txt(file):
    return StringIO(file.getvalue().decode("utf-8")).read()

def read_pptx(file):
    prs = Presentation(file)
    text_runs = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(f"[Slide {slide_num}] {shape.text}")
    return "\n".join(text_runs)

# --- Scan Functions ---
def scan_text(text, banned_dict, chars_per_page=1800):
    results, skipped = [], []
    for term, suggestions in banned_dict.items():
        pattern = re.compile(rf"\\b({re.escape(term)})\\b", re.IGNORECASE)
        for match in pattern.finditer(text):
            start = match.start()
            estimated_page = max(1, (start // chars_per_page) + 1)
            snippet = text[max(0, start - 40): min(len(text), start + 60)].replace('\n', ' ')
            if term.lower() in ["national", "taiwan"] and is_named_entity(snippet, term):
                skipped.append({"Skipped Term": match.group(), "Page": estimated_page, "Context": f"...{snippet.strip()}..."})
                continue
            results.append({
                "Banned Term": match.group(),
                "Page": estimated_page,
                "Context": f"...{snippet.strip()}...",
                "Suggested Replacement(s)": ", ".join(suggestions)
            })
    return results, skipped, text

def scan_pdf(file, banned_dict):
    results, skipped = [], []
    doc = fitz.open(stream=file.read(), filetype="pdf")
    all_text = ""
    for page_number, page in enumerate(doc, start=1):
        text = page.get_text()
        all_text += text + "\n"
        for term, suggestions in banned_dict.items():
            pattern = re.compile(rf"\\b({re.escape(term)})\\b", re.IGNORECASE)
            for match in pattern.finditer(text):
                snippet = text[max(0, match.start() - 40): min(len(text), match.end() + 60)].replace('\n', ' ')
                if term.lower() in ["national", "taiwan"] and is_named_entity(snippet, term):
                    skipped.append({"Skipped Term": match.group(), "Page": page_number, "Context": f"...{snippet.strip()}..."})
                    continue
                results.append({
                    "Banned Term": match.group(),
                    "Page": page_number,
                    "Context": f"...{snippet.strip()}...",
                    "Suggested Replacement(s)": ", ".join(suggestions)
                })
    return results, skipped, all_text

# --- Highlighting Function ---
def highlight_terms(text, terms):
    for term in sorted(terms, key=len, reverse=True):
        pattern = re.compile(rf'\\b({re.escape(term)})\\b', flags=re.IGNORECASE)
        text = pattern.sub(r"<mark style='background-color: yellow'>\\1</mark>", text)
    return text

# --- Tabs ---
tab1, tab2 = st.tabs(["\U0001F4C1 Upload Document", "\U0001F4D6 View Banned Terms"])

with tab1:
    uploaded_file = st.file_uploader("Upload a .pdf, .docx, .txt, or .pptx file", type=["pdf", "docx", "txt", "pptx"])
    if uploaded_file:
        with st.spinner("\U0001F50E Scanning document for flagged language..."):
            findings, skipped, raw_text = [], [], ""
            try:
                if uploaded_file.name.endswith(".pdf"):
                    findings, skipped, raw_text = scan_pdf(uploaded_file, banned_terms_dict)
                elif uploaded_file.name.endswith(".docx"):
                    raw_text = read_docx(uploaded_file)
                    findings, skipped, raw_text = scan_text(raw_text, banned_terms_dict)
                elif uploaded_file.name.endswith(".txt"):
                    raw_text = read_txt(uploaded_file)
                    findings, skipped, raw_text = scan_text(raw_text, banned_terms_dict)
                elif uploaded_file.name.endswith(".pptx"):
                    raw_text = read_pptx(uploaded_file)
                    findings, skipped, raw_text = scan_text(raw_text, banned_terms_dict)
            except Exception as e:
                st.error(f"Error reading file: {e}")

if 'findings' in locals() and findings is not None:
    df = pd.DataFrame(findings)

    with st.expander("📊 Summary Statistics", expanded=True):
        if not df.empty:
            st.metric("Total Banned Terms Flagged", len(df))
            st.metric("Unique Terms Found", df['Banned Term'].nunique())
            st.dataframe(df, use_container_width=True)

            term_counts = df['Banned Term'].value_counts().reset_index()
            term_counts.columns = ['Term', 'Count']
            st.subheader("📈 Most Frequently Flagged Terms")
            st.bar_chart(term_counts.set_index("Term"))
        else:
            st.warning("No banned terms were found in the uploaded document.")

    if raw_text and not df.empty:
        with st.expander("🖍️ Highlighted Text Preview", expanded=True):
            highlighted_text = highlight_terms(raw_text, df["Banned Term"].unique())
            st.markdown(
                f"<div style='white-space: pre-wrap'>{highlighted_text}</div>",
                unsafe_allow_html=True
            )
    df = pd.DataFrame(findings)

    with st.expander("📊 Summary Statistics", expanded=True):
        if not df.empty:
            st.metric("Total Banned Terms Flagged", len(df))
            st.metric("Unique Terms Found", df['Banned Term'].nunique())
            st.dataframe(df, use_container_width=True)

            term_counts = df['Banned Term'].value_counts().reset_index()
            term_counts.columns = ['Term', 'Count']
            st.subheader("📈 Most Frequently Flagged Terms")
            st.bar_chart(term_counts.set_index("Term"))
        else:
            st.warning("No banned terms were found in the uploaded document.")

    if raw_text and not df.empty:
        with st.expander("🖍 Highlighted Text Preview", expanded=True):
            highlighted_text = highlight_terms(raw_text, df["Banned Term"].unique())
            st.markdown(f"<div style='white-space: pre-wrap'>{highlighted_text}</div>", unsafe_allow_html=True)
            highlighted_text = highlight_terms(raw_text, df["Banned Term"].unique())
            st.markdown(f"<div style='white-space: pre-wrap'>{highlighted_text}</div>", unsafe_allow_html=True)

        if skipped:
            with st.expander("\u2705 Skipped Terms (Named Entities in Organization Names)", expanded=False):
                st.dataframe(pd.DataFrame(skipped), use_container_width=True)

with tab2:
    st.subheader("\U0001F4D3 Banned Terms and Suggested Replacements")
    banned_df = pd.DataFrame([
        {"Banned Term": term, "Suggested Replacement(s)": ", ".join(suggestions)}
        for term, suggestions in banned_terms_dict.items()
    ])
    st.dataframe(banned_df, use_container_width=True)

    st.markdown("""
---
### ⚖️ Explanation of Skipped Terms

Some terms like **"Taiwan"** and **"national"** are *only* skipped when they appear as part of official organization names (e.g., universities, government bodies).

#### ✅ Allowed (Skipped due to Organization Name)
- *National **Taiwan** University is a top academic institution.*
- *The **Taiwan** External Trade Development Council supports exports.*
- *The **National** Development Council issued a new report.*

#### ❌ Disallowed (Should Be Flagged)
- *“**Taiwan** is a **country** in East Asia.”*
- *“U.S. policy toward **Taiwan** has shifted.”*
- *“**National** identity is central to the reform agenda.”*
""")
